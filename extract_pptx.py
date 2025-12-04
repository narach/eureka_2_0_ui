#!/usr/bin/env python3
"""
Script to extract text from a PowerPoint (.pptx) file
"""
import sys
import os

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is not installed. Installing...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extract all text from a PowerPoint presentation"""
    prs = Presentation(file_path)
    text_content = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_content.append(f"\n{'='*60}")
        text_content.append(f"Slide {i}")
        text_content.append(f"{'='*60}\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
                text_content.append("")
    
    return "\n".join(text_content)

if __name__ == "__main__":
    pptx_file = "Eureka Concept_v1_2.pptx"
    
    if not os.path.exists(pptx_file):
        print(f"Error: {pptx_file} not found")
        sys.exit(1)
    
    try:
        text = extract_text_from_pptx(pptx_file)
        output_file = "project_description.txt"
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"Text extracted successfully to {output_file}")
        print("\n" + "="*60)
        print("EXTRACTED CONTENT:")
        print("="*60 + "\n")
        print(text)
    except Exception as e:
        import traceback
        print(f"Error extracting text: {e}")
        traceback.print_exc()
        sys.exit(1)

